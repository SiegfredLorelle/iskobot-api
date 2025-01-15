from pptx import Presentation
from typing import Dict, Any
from ..utils.text_cleaner import clean_text

def extract_text_from_pptx(blob) -> Dict[str, Any]:
    """Extract text from PPTX with metadata."""
    print(f"Processing: {blob.name}")
    metadata = {
        "source": blob.name,
        "slide_numbers": []
    }
    
    with blob.open("rb") as f:
        presentation = Presentation(f)
        text_by_slide = []
        
        for slide_num, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text.append(clean_text(paragraph.text))
            
            if slide_text:
                text_by_slide.append("\n".join(slide_text))
                metadata["slide_numbers"].append(slide_num)
        
        metadata["total_slides"] = len(presentation.slides)
    
    return {
        "text": "\n\n".join(text_by_slide),
        "metadata": metadata
    }
